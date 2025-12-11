# Licensed to the Apache Software Foundation (ASF) under one
# or more contributor license agreements.  See the NOTICE file
# distributed with this work for additional information
# regarding copyright ownership.  The ASF licenses this file
# to you under the Apache License, Version 2.0 (the
# "License"); you may not use this file except in compliance
# with the License.  You may obtain a copy of the License at
#
#   http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing,
# software distributed under the License is distributed on an
# "AS IS" BASIS, WITHOUT WARRANTIES OR CONDITIONS OF ANY
# KIND, either express or implied.  See the License for the
# specific language governing permissions and limitations
# under the License.

from typing import Union
import base64
import io

from .utils import get_plpy

def parse_pdf(file_content: Union[str, bytes]):
    """
    Parse PDF file content and extract text
    
    Extracts text from a base64-encoded PDF file or bytes. This function decodes the base64 content,
    uses PyPDF2 to extract text from each page, and concatenates the results.
    
    Args:
        file_content (Union[str, bytes]): Base64-encoded PDF file content or bytes
        
    Returns:
        str: Extracted text from the PDF file
        
    Raises:
        Warning: If input is not a string
        Warning: If the base64 decoding fails
        Warning: If the file is not a valid PDF
        Warning: If any error occurs during PDF processing
        
    Example:
        pdf_base64 = "base64_encoded_pdf_content"
        text = parse_pdf(pdf_base64)
        # Returns the extracted text from the PDF
    """
    plpy = get_plpy()
    
    try:
        import PyPDF2

        if isinstance(file_content, str):
            if file_content.startswith("data:application/pdf;base64,"):
                file_content = file_content.split(",")[1]
            file_content = base64.b64decode(file_content)
        
        pdf_file = io.BytesIO(file_content)
        
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        pages = []
        
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            pages.append({
                "page_num": page_num,
                "text": page.extract_text() or ""
            })

        return pages
    except base64.binascii.Error:
        plpy.warning("Failed to decode base64 content. Please check the input format.")
        return {
            "error": "Failed to decode base64 content. Please check the input format."
        }
    except PyPDF2.errors.PdfReadError:
        plpy.warning("Failed to read PDF file. The content may not be a valid PDF.")
        return {
            "error": "Failed to read PDF file. The content may not be a valid PDF."
        }
    except Exception as e:
        plpy.warning(f"Error processing PDF file: {e}")
        return {
            "error": f"Error processing PDF file: {e}"
        }


def parse_docx(file_content: Union[str, bytes]):
    """
    Parse DOCX file content and extract text
    
    Extracts text from a DOCX file or bytes. This function decodes the base64 content,
    uses python-docx to extract text from the document, and returns the concatenated text.
    
    Args:
        file_content (Union[str, bytes]): DOCX file content or bytes
        
    Returns:
        str: Extracted text from the DOCX file
        
    Raises:
        Warning: If input is not a string
        Warning: If the base64 decoding fails
        Warning: If the file is not a valid DOCX
        Warning: If any error occurs during DOCX processing
        
    Example:
        docx_base64 = "base64_encoded_docx_content"
        text = parse_docx(docx_base64)
        # Returns the extracted text from the DOCX
    """
    plpy = get_plpy()
    
    try:
        from docx import Document
        
        if isinstance(file_content, str):
            if file_content.startswith("data:application/vnd.openxmlformats-officedocument.wordprocessingml.document;base64,"):
                file_content = file_content.split(",")[1]
            file_content = base64.b64decode(file_content)
        
        docx_file = io.BytesIO(file_content)
        
        doc = Document(docx_file)
        paragraphs = []

        for para_num, para in enumerate(doc.paragraphs):
            paragraphs.append({
                "para_num": para_num + 1,
                "text": para.text
            })
        
        tables = []
        # Extract text from tables
        for table_num, table in enumerate(doc.tables):
            # Store table as a list of rows, each row being a list of cell texts
            table_data = []
            for row_num, row in enumerate(table.rows):
                row_cells = []
                for cell_num, cell in enumerate(row.cells):
                    row_cells.append({
                        "cell_num": cell_num + 1,
                        "text": cell.text
                    })
                table_data.append({
                    "row_num": row_num + 1,
                    "cells": row_cells
                })
            
            # Generate a plain text representation of the table for backward compatibility
            plain_text_rows = []
            for row in table_data:
                row_text = "\t".join(cell["text"] for cell in row["cells"])
                plain_text_rows.append(row_text)
            
            tables.append({
                "table_num": table_num + 1,
                "data": table_data,  # Structured table data with rows and columns
                "text": "\n".join(plain_text_rows),  # Tab-separated plain text
                "num_rows": len(table_data),
                "num_columns": len(table_data[0]["cells"]) if table_data else 0
            })
        
        return {
            "paragraphs": paragraphs,
            "tables": tables
        }
    except base64.binascii.Error:
        plpy.warning("Failed to decode base64 content. Please check the input format.")
        return {
            "error": "Failed to decode base64 content. Please check the input format."
        }
    except Exception as e:
        plpy.warning(f"Error processing DOCX file: {e}")
        return {
            "error": f"Error processing DOCX file: {e}"
        }


def parse_pptx(file_content: Union[str, bytes]):
    """
    Parse PPTX file content and extract text
    
    Extracts text from a PPTX file or bytes. This function decodes the base64 content,
    uses python-pptx to extract text from slides, and returns the concatenated text.
    
    Args:
        file_content (Union[str, bytes]): PPTX file content or bytes
        
    Returns:
        dict: Extracted text from the PPTX file, with slides as a list of strings
        
    Raises:
        Warning: If input is not a string or bytes
        Warning: If the base64 decoding fails
        Warning: If the file is not a valid PPTX
        Warning: If any error occurs during PPTX processing
        
    Example:
        pptx_base64 = "base64_encoded_pptx_content"
        text = parse_pptx(pptx_base64)
        # Returns the extracted text from the PPTX
    """
    plpy = get_plpy()
    
    try:   
        from pptx import Presentation

        if isinstance(file_content, str):
            if file_content.startswith("data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,"):
                file_content = file_content.split(",")[1]
            file_content = base64.b64decode(file_content)
        
        pptx_file = io.BytesIO(file_content)
        prs = Presentation(pptx_file)

        slides = []        
        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    slide_text.append(shape.text)
            
            if len(slide_text) > 1:
                slides.append({
                    "slide_idx": slide_idx,
                    "text": "\n".join(slide_text)
                }
            )
        
        return slides

    except base64.binascii.Error:
        plpy.warning("Failed to decode base64 content. Please check the input format.")
        return {
            "error": "Failed to decode base64 content. Please check the input format."
        }
    except Exception as e:
        plpy.warning(f"Error processing PPTX file: {e}")
        return {
            "error": f"Error processing PPTX file: {e}"
        }

def parse_document(file_content: Union[str, bytes], file_extension: str):
    """
    Universal document parser that handles different file formats

    Automatically detects the file type based on extension and calls the appropriate parser function.

    Args:
        file_content (Union[str, bytes]): File content or bytes
        file_extension (str): File extension (pdf, docx, pptx, png, etc.)
        provider_id (str, optional): Provider ID for AI-based processing (required for some formats)
        model (str, optional): Model name for AI-based processing (required for some formats)

    Returns:
        str: Extracted text from the document

    Raises:
        Warning: If input parameters are invalid
        Warning: If the file format is not supported
        Warning: If any error occurs during document processing
        
    Example:
        doc_base64 = "base64_encoded_document_content"
        text = parse_document(doc_base64, "pdf")
        # Returns the extracted text from the document
    """
    plpy = get_plpy()

    if not isinstance(file_extension, str) or not file_extension:
        plpy.warning("Input 'file_extension' must be a non-empty string.")
        return {
            "error": "Input 'file_extension' must be a non-empty string."
        }
    
    file_extension = file_extension.lower().lstrip('.')

    if file_extension == 'pdf':
        return parse_pdf(file_content)
    elif file_extension == 'docx':
        return parse_docx(file_content)
    elif file_extension == 'pptx':
        return parse_pptx(file_content)
    else:
        plpy.warning(f"Unsupported file format: {file_extension}")
        return {
            "error": f"Unsupported file format: {file_extension}"
        }
